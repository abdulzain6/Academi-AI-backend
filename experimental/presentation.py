from pptx.slide import Slide
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx import Presentation

import copy
from threading import Thread
from tempfile import NamedTemporaryFile
import os
from typing import Any, Optional, List
from langchain.chat_models import ChatOpenAI
from langchain.output_parsers import PydanticOutputParser
from exceptions import NoValidSequenceException
from langchain.prompts import (
    ChatPromptTemplate,
    SystemMessagePromptTemplate,
    HumanMessagePromptTemplate,
)
from langchain.chains import LLMChain
from langchain.chat_models import ChatOpenAI
from langchain.output_parsers import PydanticOutputParser
from database import TemplateDBManager, KnowledgeManager, TemplateModel, PlaceholderModel, SlideModel, initialize_managers
from pydantic import BaseModel, Field
from image_gen import PexelsImageSearch

class PresentationInput(BaseModel):
    topic: str
    instructions: str
    number_of_pages: int
    negative_prompt: str


class PresentationSequencePart(BaseModel):
    page_number: int = Field(
        json_schema_extra={
            "description": "The page number of the slide. Make sure this isn't repeated"
        }
    )
    slide_detail: str = Field(
        json_schema_extra={"description": "What the slide should be about."}
    )
    slide_type: str = Field(json_schema_extra={"description": "The type of the slide."})


class PlaceholderData(BaseModel):
    placeholder_name: str = Field(
        json_schema_extra={
            "description": "The name of the placeholder. Make sure name is exactly same"
        }
    )
    placeholder_data: str = Field(
        json_schema_extra={
            "description": "The data to display in the placeholder. Make sure it fits in a presentaion slide must be short (Important)"
        }
    )


class Placeholders(BaseModel):
    placeholders: Optional[List[PlaceholderData]] = Field(
        json_schema_extra={"description": "The placeholder data"}
    )


class PresentationSequence(BaseModel):
    slide_sequence: List[PresentationSequencePart] = Field(
        json_schema_extra={
            "description": "The sequence of slides used to make the presentation"
        }
    )


class PresentationMaker:
    def __init__(
        self,
        template_manager: TemplateDBManager,
        knowledge_manager: KnowledgeManager,
        openai_api_key: str,
        pexel_image_gen: PexelsImageSearch
    ) -> None:
        self.template_manager = template_manager
        self.knowledge_manager = knowledge_manager
        self.openai_api_key = openai_api_key
        self.pexel_image_gen = pexel_image_gen

    def get_best_template(self, topic: str) -> TemplateModel:
        template_name = self.knowledge_manager.get_best_template(topic)
        return self.template_manager.read_template(template_name)

    def format_placeholders(self, placeholders: List[PlaceholderModel]) -> str:
        formatted_str = "Placeholders:\n"
        if not placeholders:
            formatted_str += "No placeholders."
        for placeholder in placeholders:
            formatted_str += f"  - Name: {placeholder.name}\n"
            formatted_str += f"    Description: {placeholder.description}\n"
        return formatted_str

    def format_slides(self, slides: List[SlideModel]) -> str:
        formatted_str = "Slides Information:\n"
        formatted_str += "=" * 40 + "\n"
        for slide in slides:
            formatted_str += f"Slide Type: {slide.slide_type}\n"
            formatted_str += self.format_placeholders(slide.placeholders)
            formatted_str += "=" * 40 + "\n"
        return formatted_str

    def create_sequence(
        self, template: TemplateModel, presentation_input: PresentationInput
    ) -> PresentationSequence:
        parser = PydanticOutputParser(pydantic_object=PresentationSequence)
        prompt = ChatPromptTemplate(
            messages=[
                SystemMessagePromptTemplate.from_template(
                    """
You are an AI designed to assist in creating presentations. You are to pick a sequence of slides to create a presentation
on the topic "{topic}". THe presentation will be of {pages} pages (Important). 
You must follow the following instructions:
{instructions}
==========================================
You must not do the following:
{negative_prompt}
==========================================
"""
                ),
                HumanMessagePromptTemplate.from_template(
                    """
Here are the available slides with the placeholders in them for automatic presentation generation:
==========================================
{slides}
==========================================

Lets think step by step, Looking at the slide types and the placeholders inside them 
to create a sequence of slides that can be used to create a perfect presentation on {topic} of {pages} pages.
Do not choose slide types that are not shown to you.

{format_instructions}
"""
                ),
            ],
            input_variables=[
                "topic",
                "pages",
                "instructions",
                "negative_prompt",
                "slides",
            ],
            partial_variables={"format_instructions": parser.get_format_instructions()},
        )

        chain = LLMChain(
            prompt=prompt,
            output_parser=parser,
            llm=ChatOpenAI(openai_api_key=self.openai_api_key, temperature=0),
        )
        return chain.run(
            topic=presentation_input.topic,
            pages=presentation_input.number_of_pages,
            instructions=presentation_input.instructions,
            negative_prompt=presentation_input.negative_prompt,
            slides=self.format_slides(template.slides),
        )

    def validate_slides(
        self, template: TemplateModel, slides: PresentationSequence
    ) -> bool:
        slide_types = [slide.slide_type for slide in slides.slide_sequence]
        slides_to_choose_from = [slide.slide_type for slide in template.slides]
        return all(slide in slides_to_choose_from for slide in slide_types)
    
    @staticmethod
    def reduce_points_by_word_limit(text: str, word_limit: int) -> str:
        """
        Reduce the number of text points to fit within a word limit, ensuring that only complete points are included.
        
        Parameters:
        text (str): The original text points to be reduced.
        word_limit (int): The word limit for reduction.
        
        Returns:
        str: The reduced text with fewer points, each being complete within the word limit.
        """
        points = text.split("\n")
        reduced_points = []
        current_word_count = 0
        
        for point in points:
            point_word_count = len(point.split())
            new_word_count = current_word_count + point_word_count
            
            if new_word_count <= word_limit:
                reduced_points.append(point)
                current_word_count = new_word_count
            else:
                break
        
        return "\n".join(reduced_points)

    @staticmethod
    def reduce_paragraph_by_word_limit(paragraph: str, word_limit: int) -> str:
        """
        Reduce the size of a paragraph based on a word limit, ensuring that only complete sentences are included.
        
        Parameters:
        paragraph (str): The original paragraph to be reduced.
        word_limit (int): The word limit for reduction.
        
        Returns:
        str: The reduced paragraph containing only complete sentences within the word limit.
        """
        sentences = paragraph.split('. ')
        reduced_paragraph = ""
        current_word_count = 0

        for sentence in sentences:
            sentence_word_count = len(sentence.split())
            new_word_count = current_word_count + sentence_word_count

            # Check if adding this sentence would exceed the word limit
            if new_word_count <= word_limit:
                # Add the sentence to the reduced paragraph
                reduced_paragraph += f"{sentence}. "
                # Update the current word count
                current_word_count = new_word_count
            else:
                # Stop adding sentences once the word limit is reached
                break

        return reduced_paragraph.strip()

    def auto_reduce_text_by_words(self, text: str, word_limit_points: int, word_limit_para: int) -> str:
        """
        Automatically reduce the text size based on its type (points or paragraph).
        If the text contains both, the paragraph reduction function is used.
        
        Parameters:
        text (str): The original text to be reduced.
        word_limit (int): The word limit for reduction.
        
        Returns:
        str: The reduced text.
        """
        if text.startswith("1.") or "\n1." in text:
            return self.reduce_points_by_word_limit(text, word_limit_points)
        elif "." in text:
            return self.reduce_paragraph_by_word_limit(text, word_limit_para)
        else:
            return self.reduce_paragraph_by_word_limit(text, word_limit_para)
    
    def get_slide_content(
        self,
        sequence_part: PresentationSequencePart,
        slide: SlideModel,
        presentation_input: PresentationInput,
        all_slides: list[PresentationSequencePart],
        word_limit_para: int,
        word_limit_points: int
    ) -> Placeholders:
        parser = PydanticOutputParser(pydantic_object=Placeholders)
        prompt = ChatPromptTemplate(
            messages=[
                SystemMessagePromptTemplate.from_template(
                    """
You are an AI designed to assist in automatic presentation generation.
You will generate content for a slide by looking at placeholders and return what should they be filled with.

Make sure the placeholder content is small enough to fit a slide dont go over slide length words per placeholder.
Keep in mind lists take up more space then paragraph.

Follow the following instructions:
==============
You must fill all placeholders.
If there are no placeholders return an empty list.
{instructions}
==============

Do not do the following:
==============
Do not choose slide types that are not shown to you.
{negative_prompt}
==============
"""
                ),
                HumanMessagePromptTemplate.from_template(
                    """
=====================
The presentation is about {presentation_topic}
Here are all the slides for the presentation:
{slides}
=====================


The slide you will generate content for is {slide_detail}. This is page {page_no} of the ppt
Here are the placeholders we want to fill:
=====================
{placeholders}
=====================

Lets think step by step, Looking at the placeholders and their descriptions to fill them for the slide topic {slide_detail}. Follow all rules above! Ensure it fits in a slide
{format_instructions}
Follow the damn rules, you gave 150 words for a placeholder last time that caused error so keep it short.
Lets think step by step to accomplish this. 
        """
                ),
            ],
            input_variables=[
                "placeholders",
                "slide_detail",
                "slides",
                "presentation_topic",
                "instructions",
                "negative_prompt",
                "page_no",
            ],
            partial_variables={"format_instructions": parser.get_format_instructions()},
        )
        chain = LLMChain(
            output_parser=parser,
            prompt=prompt,
            llm=ChatOpenAI(
                openai_api_key=self.openai_api_key,
                temperature=0,
                max_tokens=250,
                model="gpt-3.5-turbo",
            ),
        )
        placeholders: Placeholders = chain.run(
            placeholders=self.format_placeholders(slide.placeholders),
            slide_detail=sequence_part.slide_detail,
            slides="\n".join([slide.slide_detail for slide in all_slides]),
            presentation_topic=presentation_input.topic,
            instructions=presentation_input.instructions,
            negative_prompt=presentation_input.negative_prompt,
            page_no=sequence_part.page_number,
        )
        for placeholder in placeholders.placeholders:
            placeholder.placeholder_data = placeholder.placeholder_data.replace(
                "\n\n", "\n"
            )
            placeholder.placeholder_data = self.auto_reduce_text_by_words(placeholder.placeholder_data, word_limit_para, word_limit_points)

        return placeholders

    def remove_placeholders_and_artifacts(self, slide) -> None:
        """
        Removes empty placeholders and other artifacts from a slide.

        Args:
        slide (Slide): The slide to clean up.

        Returns:
        None
        """
        shapes_to_remove = [
            shape
            for shape in slide.shapes
            if shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER
            and (shape.has_text_frame and not shape.text)
        ]
        for shape in shapes_to_remove:
            sp = shape._element
            sp.getparent().remove(sp)

    def duplicate_slide_with_temp_img(
        self, prs: Presentation, slide_index: int, layout_index: int
    ) -> None:
        print("Starting slide duplication...")
        slide_to_copy = prs.slides[slide_index]
        slide_layout = prs.slide_layouts[layout_index]

        print(f"Number of shapes in slide to copy: {len(slide_to_copy.shapes)}")

        new_slide = prs.slides.add_slide(slide_layout)

        print(f"Number of shapes in new slide: {len(new_slide.shapes)}")

        img_dict = {}

        for shape in slide_to_copy.shapes:
            try:
                print(f"Processing shape: {shape.name}")

                if "Picture" in shape.name:
                    with NamedTemporaryFile(suffix=".jpg", delete=False) as temp_file:
                        temp_file.write(shape.image.blob)
                    img_dict[temp_file.name] = (
                        shape.left,
                        shape.top,
                        shape.width,
                        shape.height,
                    )
                else:
                    el = shape.element
                    new_el = copy.deepcopy(el)
                    new_slide.shapes._spTree.insert_element_before(new_el, "p:extLst")

                    print(f"Shape element inserted: {new_el.tag}")

            except Exception as e:
                print(f"Could not copy shape due to: {e}")

        for img_path, dims in img_dict.items():
            new_slide.shapes.add_picture(img_path, *dims)
            os.remove(img_path)

        self.remove_placeholders_and_artifacts(new_slide)
        print(
            f"Slide duplication completed. New slide has {len(new_slide.shapes)} shapes."
        )

    def create_presentation_from_sequence(
        self,
        sequence: PresentationSequence,
        template_path: str,
        template_slides: List[SlideModel],
    ) -> Presentation:
        print("Creating a new presentation from sequence...")

        # Create the slide_type to layout_index mapping dynamically
        slide_type_to_layout_index = {
            slide.slide_type: slide.page_number - 1 for slide in template_slides
        }
        print(f"Map {slide_type_to_layout_index}")

        # Step 1: Load Template Presentation
        prs = Presentation(template_path)

        # Step 2: Iterate Through the Sequence
        for slide_part in sequence.slide_sequence:
            print(f"Processing slide type: {slide_part.slide_type}")

            if slide_model := next(
                (
                    slide
                    for slide in template_slides
                    if slide.slide_type == slide_part.slide_type
                ),
                None,
            ):
                index_to_duplicate = slide_model.page_number - 1
                print(f"Index to duplicate: {index_to_duplicate}")

                # Use the mapping to get the layout index
                layout_index = slide_type_to_layout_index.get(slide_part.slide_type, 0)
                print(f"Layout index: {layout_index}")

                # Debugging: Slide Count Before Duplication
                print(f"Slide Count Before Duplication: {len(prs.slides)}")

                self.duplicate_slide_with_temp_img(
                    prs, index_to_duplicate, layout_index
                )

                # Debugging: Slide Count After Duplication
                print(f"Slide Count After Duplication: {len(prs.slides)}")

        # Step 3: Delete All Original Slides
        for i in range(len(prs.slides) - len(sequence.slide_sequence) - 1, -1, -1):
            rId = prs.slides._sldIdLst[i].rId
            prs.part.drop_rel(rId)
            del prs.slides._sldIdLst[i]

        return prs

    def fill_single_slide(self, i: int, slide_part: Any, prs: Presentation, template_slides: List[SlideModel], presentation_input: PresentationInput, sequence: PresentationSequence, word_limit_para: int, word_limit_points: int) -> None:
        print(f"Filling in placeholders for slide: {i + 1}")
        template_slide = self.get_slide_by_type(slide_part.slide_type, template_slides)
        slide_content = self.get_slide_content(slide_part, template_slide, presentation_input, sequence.slide_sequence, word_limit_para, word_limit_points)
        presentation_slide = prs.slides[i]
        self.replace_placeholders_in_single_slide(presentation_slide, slide_content)

    def fill_presentation_with_content(
        self,
        prs: Presentation,
        sequence: PresentationSequence,
        presentation_input: PresentationInput,
        template: TemplateModel
    ) -> None:
        print("Filling in placeholders...")
        threads = []
        for i, slide_part in enumerate(sequence.slide_sequence):
            thread = Thread(target=self.fill_single_slide, args=(i, slide_part, prs, template.slides, presentation_input, sequence, template.word_limit_para, template.word_limit_points))
            threads.append(thread)
            thread.start()

        # Wait for all threads to complete
        for thread in threads:
            thread.join()

        print("Completed filling in placeholders.")

    def make_presentation(self, presentation_input: PresentationInput, slide_save_path: str) -> None:
        print("Fetching best template...")
        template = self.get_best_template(presentation_input.topic)
        slide_path = self.template_manager.get_template_file(template.template_name)

        print("Creating and validating slide sequence...")
        for _ in range(3):
            sequence = self.create_sequence(template, presentation_input)
            if self.validate_slides(template, sequence):
                break
        else:
            raise NoValidSequenceException("Couldn't find the best sequence")

        print("Sorting the slide sequence...")
        sequence.slide_sequence = sorted(
            sequence.slide_sequence, key=lambda x: x.page_number if x else 0
        )
        print(sequence)
        exit()

        prs = self.create_presentation_from_sequence(
            sequence, slide_path, template.slides
        )

        self.fill_presentation_with_content(
            prs, sequence, presentation_input, template
        )

        print("Saving the new presentation...")
        prs.save(slide_save_path)
        print("Presentation saved successfully.")

    def replace_placeholders_in_single_slide(
        self, slide: Slide, placeholders: Placeholders
    ):
        for shape in slide.shapes:
            if shape.has_text_frame:
                text_frame = shape.text_frame
                for paragraph in text_frame.paragraphs:
                    for run in paragraph.runs:
                        for placeholder in placeholders.placeholders:
                            if placeholder.placeholder_name in run.text:
                                run.text = run.text.replace(
                                    "{{" + placeholder.placeholder_name + "}}",
                                    placeholder.placeholder_data,
                                )
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        for placeholder in placeholders.placeholders:
                            if placeholder.placeholder_name in cell.text:
                                cell.text = cell.text.replace(
                                    "{{" + placeholder.placeholder_name + "}}",
                                    placeholder.placeholder_data,
                                )
        return slide

    def get_slide_by_type(self, type: str, slides: list[SlideModel]) -> SlideModel:
        for slide in slides:
            if slide.slide_type.lower() == type.lower():
                return slide
            
            
if __name__ == "__main__":
    OPENAI_API_KEY = "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"
    PEXELS_API_KEY = "rX8ysruEFR2U4IMpCjh9KviIdKl0orDwJRXwmf6mVRkbEGlbdyATveM8"
    
    import langchain
    langchain.verbose = True
    
    template_manager, knowledge_manager = initialize_managers(OPENAI_API_KEY)
    presentation_maker = PresentationMaker(
        template_manager, knowledge_manager, OPENAI_API_KEY, PexelsImageSearch(PEXELS_API_KEY)
    )
    
    
    
    
    
    print(
        presentation_maker.make_presentation(
            PresentationInput(
                topic="science project",
                instructions="Explain as if i was 10",
                number_of_pages=5,
                negative_prompt="use hard vocabulary",
            ),
            "example_modified.pptx"
        )
    )
