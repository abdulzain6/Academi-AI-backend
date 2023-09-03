from pptx.util import Inches
import firebase
import os
from shutil import copyfile
from pydantic import BaseModel
from typing import Optional, List, Union
from firebase_admin import firestore, storage
from abc import ABC, abstractmethod
from langchain.vectorstores import Qdrant
from langchain.embeddings import OpenAIEmbeddings
from langchain.schema import Document
from qdrant_client import QdrantClient
from qdrant_client.http import models as rest
from pptx import Presentation
from typing import Union


class PlaceholderModel(BaseModel):
    name: str
    description: Optional[str] = None
    is_image: bool = False
    image_width: Optional[int] = None
    image_height: Optional[int] = None


class SlideModel(BaseModel):
    slide_type: str
    page_number: int
    placeholders: List[PlaceholderModel]


class TemplateModel(BaseModel):
    template_name: str
    template_description: str
    category: Optional[str] = None
    version: Optional[str] = None
    aspect_ratio: Optional[str] = None
    color_scheme: Optional[List[str]] = None
    slides: List[SlideModel]
    file_extension: str = ".pptx"
    word_limit_para: int = 93
    word_limit_points: int = 80
    word_limit_hybrid: int = 75



class TemplateObserver(ABC):
    @abstractmethod
    def update(self, template: TemplateModel) -> None:
        pass


class TemplateListManager(TemplateObserver):
    def __init__(self) -> None:
        self.templates: List[TemplateModel] = []

    def update(self, template: TemplateModel) -> None:
        self.templates.append(template)


class TemplateDBManager:
    def __init__(self, local_storage_path: str = "temp") -> None:
        self.db = firestore.client()
        self.template_collection = self.db.collection("templates")
        self.bucket = storage.bucket()
        self.local_storage_path = local_storage_path
        self._observers: List[TemplateObserver] = []
        if not os.path.exists(self.local_storage_path):
            os.makedirs(self.local_storage_path)

    def register_observer(self, observer: TemplateObserver) -> None:
        self._observers.append(observer)

    def _notify_observers(self, template: TemplateModel) -> None:
        for observer in self._observers:
            observer.update(template)

    def is_template_name_unique(self, template_name: str) -> bool:
        return self._find_template_by_name(template_name) is None

    def _find_template_by_name(self, template_name: str):
        """Helper function to find a template by its name."""
        docs = self.template_collection.where(
            "template_name", "==", template_name
        ).stream()
        for doc in docs:
            return doc
        return None

    def create_template(
        self, template: TemplateModel, file_path: str
    ) -> Union[None, str]:
        if not self.is_template_name_unique(template.template_name):
            return "Template name already exists."

        file_extension = os.path.splitext(file_path)[1]
        template.file_extension = file_extension

        template_dict = template.model_dump()
        _, doc_ref = self.template_collection.add(template_dict)

        doc_id = doc_ref.id
        self._notify_observers(template)

        self.upload_template_file(file_path, doc_id, file_extension)

        return None

    def read_template(self, template_name: str) -> Union[TemplateModel, None]:
        doc = self._find_template_by_name(template_name)
        return TemplateModel(**doc.to_dict()) if doc else None

    def update_template(self, template_name: str, updated_data: dict) -> None:
        if doc := self._find_template_by_name(template_name):
            doc.reference.update(updated_data)

    def delete_template(self, template_name: str) -> None:
        if doc := self._find_template_by_name(template_name):
            doc_data = doc.to_dict()
            file_extension = doc_data.get("file_extension", "")
            doc_id = doc.id
            doc.reference.delete()
            self.delete_template_file(doc_id, file_extension)

    def delete_template_file(self, doc_id: str, file_extension: str) -> None:
        blob = self.bucket.blob(f"templates/{doc_id}{file_extension}")
        blob.delete()
        local_path = os.path.join(self.local_storage_path, f"{doc_id}{file_extension}")
        if os.path.exists(local_path):
            os.remove(local_path)

    def get_all_templates(self) -> List[TemplateModel]:
        docs = self.template_collection.stream()
        return [TemplateModel(**doc.to_dict()) for doc in docs]

    def get_template_file(self, template_name: str) -> Union[str, None]:
        doc = self._find_template_by_name(
            template_name
        )  # Assuming _find_template_by_name returns the Firestore document
        if doc is None:
            return None

        doc_data = doc.to_dict()
        file_extension = doc_data.get(
            "file_extension", ".pptx"
        )  # Fetch the file extension from the document
        doc_id = doc.id

        filename = f"{doc_id}{file_extension}"
        local_path = os.path.join(self.local_storage_path, filename)

        if os.path.exists(local_path):
            return local_path

        cloud_path = f"templates/{filename}"
        blob = self.bucket.blob(cloud_path)
        blob.download_to_filename(local_path)

        return local_path

    def upload_template_file(
        self, file_path: str, doc_id: str, file_extension: str
    ) -> None:
        blob = self.bucket.blob(f"templates/{doc_id}{file_extension}")
        blob.upload_from_filename(file_path)
        local_path = os.path.join(self.local_storage_path, f"{doc_id}{file_extension}")
        copyfile(file_path, local_path)


class KnowledgeManager(TemplateObserver):
    def __init__(self, openai_api_key: str) -> None:
        self.openai_api_key = openai_api_key
        self.vectorstore = self.get_vectorstore()

    def get_vectorstore(self) -> Qdrant:
        client = QdrantClient(location=":memory:")
        embedding = OpenAIEmbeddings(openai_api_key=self.openai_api_key)
        partial_embeddings = embedding.embed_documents(["test"])
        vector_size = len(partial_embeddings[0])

        vectors_config = rest.VectorParams(
            size=vector_size,
            distance=rest.Distance.COSINE,
        )
        client.create_collection("templates", vectors_config=vectors_config)
        return Qdrant(embeddings=embedding, client=client, collection_name="templates")

    def add_documents(self, documents: List[Document]) -> None:
        self.vectorstore.add_documents(documents)

    def update(self, template: TemplateModel) -> None:
        document = Document(
            page_content=f"{template.template_description} {template.template_name}",
            metadata={"name": template.template_name},
        )
        self.add_documents([document])

    def get_best_template(self, topic: str) -> str:
        return self.vectorstore.similarity_search(topic, k=1)[0].metadata["name"]


def initialize_managers(
    openai_api_key: str, local_storage_path: str = "temp"
) -> tuple[TemplateDBManager, KnowledgeManager]:
    template_manager = TemplateDBManager(local_storage_path)
    knowledge_manager = KnowledgeManager(openai_api_key)
    templates = template_manager.get_all_templates()
    docs = [
        Document(
            page_content=f"{template.template_description} {template.template_name}",
            metadata={"name": template.template_name},
        )
        for template in templates
    ]
    knowledge_manager.add_documents(docs)
    template_manager.register_observer(knowledge_manager)
    return template_manager, knowledge_manager


class PresentationTemplateWizard:
    def __init__(self, ppt_file_path: str):
        self.ppt = Presentation(ppt_file_path)

    def get_image_dimensions(
        self, slide_number: int, placeholder_name: str
    ) -> Union[None, tuple]:
        slide = self.ppt.slides[slide_number - 1]  # Adjust for 1-based indexing
        
        for shape in slide.shapes:
            if shape.shape_type == 13:  # Shape type 13 refers to a picture
                if shape.name == placeholder_name:
                    width_px = int(shape.width / Inches(1) * 96)  # Conversion to pixels
                    height_px = int(shape.height / Inches(1) * 96)  # Conversion to pixels
                    return width_px, height_px

        return None  # Return None if placeholder not found

    def wizard_cli(self) -> TemplateModel:
        print("Welcome to the Presentation Template Wizard!")

        template_name = input("Enter the template name: ")
        template_description = input("Enter a description for the template: ")
        category = input("Enter the category (optional, press Enter to skip): ") or None
        version = input("Enter the version (optional, press Enter to skip): ") or None
        aspect_ratio = (
            input("Enter the aspect ratio (optional, press Enter to skip): ") or None
        )
        color_scheme = (
            input(
                "Enter color scheme as comma-separated values (optional, press Enter to skip): "
            ).split(",")
            if input("Do you have a color scheme? (y/n): ").lower() == "y"
            else None
        )
        word_limit_para = int(input("Enter the word limit for paragraphs: "))
        word_limit_points = int(input("Enter the word limit for points: "))

        slides = []

        while input("Do you want to add a slide? (y/n): ").lower() == "y":
            slide_type = input("Enter the slide type: ")
            page_number = int(input("Enter the page number: "))
            placeholders = []

            while input("Do you want to add a placeholder? (y/n): ").lower() == "y":
                name = input("Enter the name of the placeholder: ")
                is_image = input("Is this an image placeholder? (y/n): ").lower() == "y"

                description = (
                    "Write a Google search query for this image that is relevant to the topic."
                    if is_image
                    else input("Enter a description (optional, press Enter to skip): ")
                    or None
                )

                if is_image:
                    if dimensions := self.get_image_dimensions(page_number, name):
                        image_width, image_height = dimensions
                    else:
                        print(
                            "Couldn't find the placeholder in the slide, skipping image dimensions."
                        )
                        image_width, image_height = None, None
                else:
                    image_width, image_height = None, None

                placeholders.append(
                    PlaceholderModel(
                        name=name,
                        description=description,
                        is_image=is_image,
                        image_width=image_width,
                        image_height=image_height,
                    )
                )

            slides.append(
                SlideModel(
                    slide_type=slide_type,
                    page_number=page_number,
                    placeholders=placeholders,
                )
            )

        return TemplateModel(
            template_name=template_name,
            template_description=template_description,
            category=category,
            version=version,
            aspect_ratio=aspect_ratio,
            color_scheme=color_scheme,
            slides=slides,
            word_limit_para=word_limit_para,
            word_limit_points=word_limit_points,
        )


if __name__ == "__main__":
    OPENAI_API_KEY = "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"    
    template_manager, knowledge_manager = initialize_managers(OPENAI_API_KEY)
    wizard = PresentationTemplateWizard("/home/zain/Downloads/minimal beautiful.pptx")
    template_model = wizard.wizard_cli()
    template_manager.create_template(template_model, "/home/zain/Downloads/minimal beautiful.pptx")