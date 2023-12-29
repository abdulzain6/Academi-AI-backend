import json
import os
from api.lib.anyscale_embeddings import AnyscaleEmbeddings
from pptx.util import Inches
from shutil import copy
from pydantic import BaseModel
from typing import Optional, List, Union
from abc import ABC, abstractmethod
from langchain.vectorstores import Qdrant
from langchain.embeddings.base import Embeddings
from langchain.schema import Document
from qdrant_client import QdrantClient
from qdrant_client.http import models as rest
from pptx import Presentation
from typing import List, Union, Optional
from langchain.embeddings.openai import OpenAIEmbeddings

DEFAULT_TEMPLATE_DIR = os.path.join(os.path.dirname(os.path.realpath(__file__)), "template_dir")
DEFAULT_TEMPLATES_JSON = os.path.join(os.path.dirname(os.path.realpath(__file__)), "templates.json")

class PlaceholderModel(BaseModel):
    name: str
    description: Optional[str] = None
    is_image: bool = False
    image_width: Optional[int] = None
    image_height: Optional[int] = None


class SlideModel(BaseModel):
    slide_type: str
    page_number: int
    placeholders: List[PlaceholderModel]


class TemplateModel(BaseModel):
    template_name: str
    template_description: str
    category: Optional[str] = None
    version: Optional[str] = None
    aspect_ratio: Optional[str] = None
    color_scheme: Optional[List[str]] = None
    slides: List[SlideModel]
    file_extension: str = "pptx"
    word_limit_para: int = 93
    word_limit_points: int = 80
    word_limit_hybrid: int = 75
    file_name: Optional[str] = None 
    image_base64: Optional[str] = None


class TemplateObserver(ABC):
    @abstractmethod
    def update(self, template: TemplateModel) -> None:
        pass


class TemplateListManager(TemplateObserver):
    def __init__(self) -> None:
        self.templates: List[TemplateModel] = []

    def update(self, template: TemplateModel) -> None:
        self.templates.append(template)


class TemplateDBManager:
    def __init__(
        self,
        templates_json_path: str,
        template_dir: str,
    ) -> None:
        self.template_dir = template_dir
        self.templates_json_path = templates_json_path
        self.templates = self.load_templates(templates_json_path)
        self._observers: List[TemplateObserver] = []
        
    def add_template_to_json(self, template_json_file: str, template: TemplateModel):
        with open(template_json_file, "r+") as fp:
            self.templates.append(template)
            json.dump(self.templates, fp)
        
    def load_templates(self, template_json_file: str = "templates.json") -> List[TemplateModel]:
        with open(template_json_file, "rt") as fp:
            templates = json.load(fp)

        return [TemplateModel.model_validate(template) for template in templates]

    def register_observer(self, observer: TemplateObserver) -> None:
        self._observers.append(observer)

    def _notify_observers(self, template: TemplateModel) -> None:
        for observer in self._observers:
            observer.update(template)

    def find_template_by_name(self, template_name: str) -> Optional[TemplateModel]:
        for template in self.templates:
            if template.template_name == template_name:
                return template

    def create_template(
        self, template: TemplateModel, file_path: str
    ) -> Union[None, str]:
        if not self.find_template_by_name(template.template_name):
            raise ValueError("Template Already exists")

        copy(file_path, self.template_dir)
        self.add_template_to_json(self.templates_json_path, template)
        self._notify_observers(template)
        return None

    def read_template(self, template_name: str) -> Union[TemplateModel, None]:
        temp = self.find_template_by_name(template_name)
        return temp if temp else None

    def get_all_templates(self) -> List[TemplateModel]:
        return self.templates

    def get_template_file(self, template_name: str) -> Union[str, None]:
        doc = self.find_template_by_name(template_name)
        return os.path.join(self.template_dir, doc.file_name)


class TemplateKnowledgeManager(TemplateObserver):
    def __init__(self, embeddings = None) -> None:
        if not embeddings:
            self.embeddings = AnyscaleEmbeddings(
                base_url="https://api.endpoints.anyscale.com/v1",
                model="thenlper/gte-large",
               max_retries=2,
               timeout=5
            )
        else:
            self.embeddings = embeddings
        try:
            self.vectorstore = self.get_vectorstore()
        except Exception:
            self.embeddings = OpenAIEmbeddings()
            self.vectorstore = self.get_vectorstore()

    def get_vectorstore(self) -> Qdrant:
        client = QdrantClient(location=":memory:")
        embedding = self.embeddings
        partial_embeddings = embedding.embed_documents(["test"])
        vector_size = len(partial_embeddings[0])

        vectors_config = rest.VectorParams(
            size=vector_size,
            distance=rest.Distance.COSINE,
        )
        client.create_collection("templates", vectors_config=vectors_config)
        return Qdrant(embeddings=embedding, client=client, collection_name="templates")

    def add_documents(self, documents: List[Document]) -> None:
        self.vectorstore.add_documents(documents)

    def update(self, template: TemplateModel) -> None:
        document = Document(
            page_content=f"{template.template_description} {template.template_name}",
            metadata={"name": template.template_name},
        )
        self.add_documents([document])

    def get_best_template(self, topic: str) -> str:
        return self.vectorstore.similarity_search(topic, k=1)[0].metadata["name"]


def initialize_managers(
    template_json_path: str, template_dir: str
) -> tuple[TemplateDBManager, TemplateKnowledgeManager]:
    template_manager = TemplateDBManager(templates_json_path=template_json_path, template_dir=template_dir)
    knowledge_manager = TemplateKnowledgeManager()
    templates = template_manager.get_all_templates()
    docs = [
        Document(
            page_content=f"{template.template_description} {template.template_name}",
            metadata={"name": template.template_name},
        )
        for template in templates
    ]
    knowledge_manager.add_documents(docs)
    template_manager.register_observer(knowledge_manager)
    return template_manager, knowledge_manager


class PresentationTemplateWizard:
    def __init__(self, ppt_file_path: str):
        self.ppt = Presentation(ppt_file_path)

    def get_image_dimensions(
        self, slide_number: int, placeholder_name: str
    ) -> Union[None, tuple]:
        slide = self.ppt.slides[slide_number - 1]  # Adjust for 1-based indexing

        for shape in slide.shapes:
            if shape.name == placeholder_name and shape.shape_type == 13:
                width_px = int(shape.width / Inches(1) * 96)  # Conversion to pixels
                height_px = int(shape.height / Inches(1) * 96)  # Conversion to pixels
                return width_px, height_px

        return None  # Return None if placeholder not found

    def wizard_cli(self) -> TemplateModel:
        print("Welcome to the Presentation Template Wizard!")

        template_name = input("Enter the template name: ")
        file_name = input("Enter filename ")
        template_description = input("Enter a description for the template: ")
        category = input("Enter the category (optional, press Enter to skip): ") or None
        version = input("Enter the version (optional, press Enter to skip): ") or None
        aspect_ratio = (
            input("Enter the aspect ratio (optional, press Enter to skip): ") or None
        )
        color_scheme = (
            input(
                "Enter color scheme as comma-separated values (optional, press Enter to skip): "
            ).split(",")
            if input("Do you have a color scheme? (y/n): ").lower() == "y"
            else None
        )
        word_limit_para = int(input("Enter the word limit for paragraphs: "))
        word_limit_points = int(input("Enter the word limit for points: "))

        slides = []

        while input("Do you want to add a slide? (y/n): ").lower() == "y":
            slide_type = input("Enter the slide type: ")
            page_number = int(input("Enter the page number: "))
            placeholders = []

            while input("Do you want to add a placeholder? (y/n): ").lower() == "y":
                name = input("Enter the name of the placeholder: ")
                is_image = input("Is this an image placeholder? (y/n): ").lower() == "y"

                description = (
                    "Write a Google search query for this image that is relevant to the topic."
                    if is_image
                    else input("Enter a description (optional, press Enter to skip): ")
                    or None
                )

                if is_image:
                    if dimensions := self.get_image_dimensions(page_number, name):
                        image_width, image_height = dimensions
                    else:
                        print(
                            "Couldn't find the placeholder in the slide, skipping image dimensions."
                        )
                        image_width, image_height = None, None
                else:
                    image_width, image_height = None, None

                placeholders.append(
                    PlaceholderModel(
                        name=name,
                        description=description,
                        is_image=is_image,
                        image_width=image_width,
                        image_height=image_height,
                    )
                )

            slides.append(
                SlideModel(
                    slide_type=slide_type,
                    page_number=page_number,
                    placeholders=placeholders,
                )
            )

        return TemplateModel(
            template_name=template_name,
            template_description=template_description,
            category=category,
            version=version,
            aspect_ratio=aspect_ratio,
            color_scheme=color_scheme,
            slides=slides,
            word_limit_para=word_limit_para,
            word_limit_points=word_limit_points,
            file_name=file_name
        )

    